"""
PPT Presentation Generator

Converts the generated JSON slide structure into a PowerPoint presentation
using python-pptx.  Field names are aligned to the schema defined in
prompts/presentation_slides_generator_prompt.txt — the code is NOT tuned to
any single JSON output.
"""

import json
import logging
import re
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Optional

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Logging — no print() anywhere in this module
# ---------------------------------------------------------------------------
logging.basicConfig(level=logging.INFO, format="%(levelname)-8s %(message)s")
log = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Layout constants
# ---------------------------------------------------------------------------
SLIDE_WIDTH  = Inches(10)
SLIDE_HEIGHT = Inches(7.5)
MARGIN_L     = Inches(0.5)
CONTENT_W    = Inches(9.0)
FOOTER_TOP   = Inches(7.12)
FOOTER_H     = Inches(0.25)
FOOTER_TEXT  = "InsightSphere"

HEADER_IMAGE = Path("templates") / "header.png"

# Content-slide title hard cap (user requirement: <=28 pt or title overflows)
CONTENT_TITLE_PT = 28

# Chart-type mapping from JSON schema string → python-pptx enum
CHART_TYPE_MAP = {
    "bar":     XL_CHART_TYPE.BAR_CLUSTERED,
    "column":  XL_CHART_TYPE.COLUMN_CLUSTERED,
    "line":    XL_CHART_TYPE.LINE_MARKERS,
    "pie":     XL_CHART_TYPE.PIE,
    "donut":   XL_CHART_TYPE.DOUGHNUT,
    "scatter": XL_CHART_TYPE.XY_SCATTER,
}

# ---------------------------------------------------------------------------
# Fallback palette (used when design_system_reference is incomplete)
# ---------------------------------------------------------------------------
FB_PRIMARY = "#02428E"
FB_ACCENT  = "#F26633"
FB_TEXT    = "#000000"
FB_WHITE   = "#FFFFFF"
FB_BG      = "#FFFFFF"


# ---------------------------------------------------------------------------
# DesignSystem  — maps the 10 fields from the prompt schema
# ---------------------------------------------------------------------------
@dataclass
class DesignSystem:
    """
    Maps exactly to design_system_reference in the prompt schema:
      theme_name, primary_color, accent_color, background_color,
      text_color_on_light, text_color_on_dark,
      font_title, font_body, font_size_title_pt, font_size_body_pt
    """
    theme_name:     str   = "InsightSphere"
    primary_color:  str   = FB_PRIMARY
    accent_color:   str   = FB_ACCENT
    bg_color:       str   = FB_BG
    text_on_light:  str   = FB_PRIMARY
    text_on_dark:   str   = FB_WHITE
    font_title:     str   = "Arial"
    font_body:      str   = "Arial"
    title_pt:       int   = 44
    body_pt:        int   = 15

    @classmethod
    def from_dict(cls, d: dict) -> "DesignSystem":
        return cls(
            theme_name    = d.get("theme_name",           cls.theme_name),
            primary_color = d.get("primary_color",        FB_PRIMARY),
            accent_color  = d.get("accent_color",         FB_ACCENT),
            bg_color      = d.get("background_color",     FB_BG),
            text_on_light = d.get("text_color_on_light",  d.get("primary_color", FB_PRIMARY)),
            text_on_dark  = d.get("text_color_on_dark",   FB_WHITE),
            font_title    = d.get("font_title",           "Arial"),
            font_body     = d.get("font_body",            "Arial"),
            title_pt      = int(d.get("font_size_title_pt", 44)),
            body_pt       = int(d.get("font_size_body_pt",  15)),
        )


# ---------------------------------------------------------------------------
# Colour helper
# ---------------------------------------------------------------------------
def _rgb(hex_color: str, fallback: str = FB_TEXT) -> RGBColor:
    try:
        h = hex_color.lstrip("#")
        if len(h) != 6:
            raise ValueError
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    except Exception:
        log.warning("Bad colour %r — falling back to %r", hex_color, fallback)
        return _rgb(fallback)


# ---------------------------------------------------------------------------
# Low-level shape / text helpers
# ---------------------------------------------------------------------------
def _blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _bg_solid(slide, hex_color: str) -> None:
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = _rgb(hex_color, FB_BG)


def _bg_gradient(slide, colors: List[str], angle: float = 45.0,
                 fallback: str = FB_PRIMARY) -> None:
    f = slide.background.fill
    f.gradient()
    f.gradient_angle = angle
    f.gradient_stops[0].color.rgb = _rgb(colors[0] if colors else fallback, fallback)
    f.gradient_stops[1].color.rgb = _rgb(
        colors[1] if len(colors) > 1 else fallback, fallback)


def _textbox(slide, left, top, width, height, wrap: bool = True):
    """Return (shape, text_frame, first_paragraph)."""
    s  = slide.shapes.add_textbox(left, top, width, height)
    tf = s.text_frame
    tf.word_wrap = wrap
    return s, tf, tf.paragraphs[0]


def _fmt(para, text: str, pt: int, bold=False, italic=False,
         color: str = FB_TEXT, align=PP_ALIGN.LEFT,
         font_name: str = None) -> None:
    para.text           = text
    para.font.size      = Pt(pt)
    para.font.bold      = bold
    para.font.italic    = italic
    para.font.color.rgb = _rgb(color)
    para.alignment      = align
    if font_name:
        para.font.name = font_name


def _accent_line(slide, top, left=None, width=None,
                 color: str = FB_ACCENT) -> None:
    left  = left  if left  is not None else MARGIN_L
    width = width if width is not None else CONTENT_W
    ln = slide.shapes.add_shape(1, left, top, width, Inches(0.045))
    ln.fill.solid()
    ln.fill.fore_color.rgb = _rgb(color)
    ln.line.color.rgb      = _rgb(color)


def _colored_box(slide, text: str, left, top, width, height,
                 bg: str, fg: str = FB_WHITE, pt: int = 14,
                 bold: bool = True, font_name: str = None) -> None:
    box = slide.shapes.add_shape(1, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = _rgb(bg)
    box.line.color.rgb      = _rgb(bg)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _fmt(tf.paragraphs[0], text, pt, bold=bold, color=fg,
         align=PP_ALIGN.CENTER, font_name=font_name)


# ---------------------------------------------------------------------------
# Reusable composite helpers
# ---------------------------------------------------------------------------
def _title_box(slide, text: str, ds: DesignSystem,
               top=None, left=None, width=None,
               align=PP_ALIGN.LEFT, color: str = None,
               pt: int = None) -> None:
    top   = top   if top   is not None else Inches(0.3)
    left  = left  if left  is not None else MARGIN_L
    width = width if width is not None else CONTENT_W
    color = color or ds.text_on_light
    pt    = pt    if pt    is not None else CONTENT_TITLE_PT
    _, _, p = _textbox(slide, left, top, width, Inches(0.75))
    _fmt(p, text, pt, bold=True, color=color, align=align, font_name=ds.font_title)


def _subtitle_box(slide, text: str, ds: DesignSystem,
                  top=None, left=None, width=None,
                  align=PP_ALIGN.LEFT, color: str = None,
                  pt: int = None) -> None:
    if not text:
        return
    top   = top   if top   is not None else Inches(1.1)
    left  = left  if left  is not None else MARGIN_L
    width = width if width is not None else CONTENT_W
    color = color or ds.text_on_light
    pt    = pt    if pt    is not None else 18
    _, _, p = _textbox(slide, left, top, width, Inches(0.45))
    _fmt(p, text, pt, bold=True, color=color, font_name=ds.font_body, align=align)


def _render_bullets(slide, items: List[dict], ds: DesignSystem,
                    left, top, width, height) -> None:
    """Render bullet dicts per prompt schema: {level, text, emphasis, accent}."""
    if not items:
        return
    _, tf, _ = _textbox(slide, left, top, width, height)
    for i, b in enumerate(items):
        para = tf.paragraphs[i] if i == 0 else tf.add_paragraph()
        lvl  = max(0, int(b.get("level", 1)) - 1)
        para.level     = lvl
        para.text      = b.get("text", "")
        para.font.size = Pt(ds.body_pt if lvl == 0 else ds.body_pt - 2)
        para.font.name = ds.font_body
        para.font.bold = (b.get("emphasis", "normal") == "bold")

        accent_val = b.get("accent", "none")
        if accent_val == "primary":
            para.font.color.rgb = _rgb(ds.primary_color)
        elif accent_val == "accent":
            para.font.color.rgb = _rgb(ds.accent_color)
        else:
            para.font.color.rgb = _rgb(ds.text_on_light)

        para.space_before = Pt(5)
        para.space_after  = Pt(5)


def _render_str_list(slide, items: List[str], color: str,
                     left, top, width, height, pt: int = 14,
                     font_name: str = None, bullet_prefix: str = "\u2022 ") -> None:
    """Render a plain list[str] into a text box with optional bullet prefix."""
    if not items:
        return
    _, tf, _ = _textbox(slide, left, top, width, height)
    for i, text in enumerate(items):
        para = tf.paragraphs[i] if i == 0 else tf.add_paragraph()
        para.text           = f"{bullet_prefix}{text}" if bullet_prefix else text
        para.font.size      = Pt(pt)
        para.font.color.rgb = _rgb(color)
        para.space_before   = Pt(4)
        para.space_after    = Pt(4)
        if font_name:
            para.font.name = font_name


def _footer(slide, ds: DesignSystem) -> None:
    _, _, p = _textbox(slide, Inches(0.3), FOOTER_TOP, Inches(3), FOOTER_H, wrap=False)
    _fmt(p, FOOTER_TEXT, 10, color=ds.primary_color, font_name=ds.font_body)


def _add_speaker_notes(slide, sd: dict) -> None:
    """Add prospect_relevance_note + slide_purpose as speaker notes."""
    parts: list[str] = []
    prn = sd.get("prospect_relevance_note", "")
    sp  = sd.get("slide_purpose", "")
    if sp:
        parts.append(f"[Slide Purpose] {sp}")
    if prn:
        parts.append(f"[Prospect Relevance] {prn}")
    if parts:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = "\n\n".join(parts)


# ---------------------------------------------------------------------------
# Schema-aware accessors (per the prompt schema)
# ---------------------------------------------------------------------------
def _get_bullets(sd: dict) -> List[dict]:
    """Prompt schema: bullets lives at sd['bullets']. Fallback to sd['content']['bullets']."""
    return sd.get("bullets") or sd.get("content", {}).get("bullets") or []


def _get_callout(sd: dict) -> Optional[dict]:
    """Prompt schema: callout_box at sd['callout_box']. Fallback to design_notes."""
    cb = sd.get("callout_box") or sd.get("design_notes", {}).get("callout_box")
    return cb if (cb and cb.get("text")) else None


def _accent_line_wanted(sd: dict) -> bool:
    val = sd.get("design_notes", {}).get("accent_line", False)
    return val is True or str(val).lower() in ("true", "yes")


def _get_accent_line_color(sd: dict, ds: DesignSystem) -> str:
    return sd.get("design_notes", {}).get("accent_line_color", ds.accent_color)


def _resolve_bg(slide, dn: dict, ds: DesignSystem,
                dark_default: str = None) -> None:
    """Apply background from design_notes: solid hex, 'gradient', or gradient_colors."""
    bg = dn.get("background", dark_default or ds.bg_color)
    gc = dn.get("gradient_colors", [])
    if gc:
        _bg_gradient(slide, gc, fallback=ds.primary_color)
    elif bg == "gradient":
        _bg_gradient(slide, [ds.accent_color, ds.primary_color],
                     fallback=ds.primary_color)
    else:
        _bg_solid(slide, bg)


# ---------------------------------------------------------------------------
# Slide creators  — one per slide_type in the prompt schema
# ---------------------------------------------------------------------------

# ---- title ----------------------------------------------------------------
def _slide_title(prs: Presentation, sd: dict, ds: DesignSystem) -> None:
    slide = _blank_slide(prs)
    dn = sd.get("design_notes", {})
    _resolve_bg(slide, dn, ds, dark_default=ds.primary_color)

    title_color    = dn.get("title_color",    ds.text_on_dark)
    subtitle_color = dn.get("subtitle_color", ds.text_on_dark)

    _, _, p = _textbox(slide, MARGIN_L, Inches(2.2), CONTENT_W, Inches(1.9))
    _fmt(p, sd.get("title", ""), ds.title_pt,
         bold=True, color=title_color, align=PP_ALIGN.CENTER, font_name=ds.font_title)

    _, _, p = _textbox(slide, MARGIN_L, Inches(4.2), CONTENT_W, Inches(0.9))
    _fmt(p, sd.get("subtitle", ""), 20,
         color=subtitle_color, align=PP_ALIGN.CENTER, font_name=ds.font_body)

    presenter = sd.get("presenter_info", "")
    if presenter:
        _, _, p = _textbox(slide, MARGIN_L, Inches(5.2), CONTENT_W, Inches(0.5))
        _fmt(p, presenter, 14, color=subtitle_color, align=PP_ALIGN.CENTER,
             font_name=ds.font_body)

    if HEADER_IMAGE.exists():
        try:
            slide.shapes.add_picture(
                str(HEADER_IMAGE),
                Inches(7.5), Inches(0.15),
                width=Inches(2.2), height=Inches(0.36),
            )
        except Exception as e:
            log.warning("Header image failed: %s", e)


# ---- content --------------------------------------------------------------
def _slide_content(prs: Presentation, sd: dict, ds: DesignSystem) -> None:
    slide = _blank_slide(prs)
    dn = sd.get("design_notes", {})
    _bg_solid(slide, dn.get("background", ds.bg_color))

    _title_box(slide, sd.get("title", ""), ds, top=Inches(0.3), pt=CONTENT_TITLE_PT)

    accent_color = _get_accent_line_color(sd, ds)
    if _accent_line_wanted(sd):
        _accent_line(slide, top=Inches(0.95), color=accent_color)

    subtitle    = sd.get("subtitle", "")
    has_subtitle = bool(subtitle)
    if has_subtitle:
        _subtitle_box(slide, subtitle, ds, top=Inches(1.05))

    bullets = _get_bullets(sd)
    callout = _get_callout(sd)
    b_top = Inches(1.55) if has_subtitle else Inches(1.2)
    b_h   = Inches(4.3)  if callout      else Inches(5.6)

    _render_bullets(slide, bullets, ds,
                    left=Inches(0.6), top=b_top,
                    width=Inches(8.8), height=b_h)

    if callout:
        _colored_box(
            slide, callout["text"],
            left=Inches(0.6), top=Inches(6.05),
            width=Inches(8.8), height=Inches(0.75),
            bg=callout.get("background_color", ds.accent_color),
            fg=callout.get("text_color", FB_WHITE),
            pt=14, font_name=ds.font_body,
        )


# ---- two_column -----------------------------------------------------------
def _slide_two_column(prs: Presentation, sd: dict, ds: DesignSystem) -> None:
    """
    Prompt schema: left_column / right_column, each with
    header, width_pct, content_type, items[], chart_or_image_description.
    design_notes: gutter_px, left_accent_color, right_accent_color.
    """
    slide = _blank_slide(prs)
    dn = sd.get("design_notes", {})
    _bg_solid(slide, dn.get("background", ds.bg_color))
    _title_box(slide, sd.get("title", ""), ds, top=Inches(0.3), pt=CONTENT_TITLE_PT)
    _accent_line(slide, top=Inches(0.95), color=ds.accent_color)

    col_defs = [
        (sd.get("left_column",  {}), Inches(0.5),
         dn.get("left_accent_color",  ds.accent_color)),
        (sd.get("right_column", {}), Inches(5.3),
         dn.get("right_accent_color", ds.primary_color)),
    ]
    col_w   = Inches(4.5)
    col_top = Inches(1.15)

    for col, col_left, hdr_color in col_defs:
        if not col:
            continue
        header = col.get("header", "")
        items  = col.get("items", [])

        _, _, hp = _textbox(slide, col_left, col_top, col_w, Inches(0.45))
        _fmt(hp, header, 18, bold=True, color=hdr_color, font_name=ds.font_title)

        _accent_line(slide, top=col_top + Inches(0.49),
                     left=col_left, width=col_w, color=hdr_color)

        _render_str_list(slide, items, ds.text_on_light,
                         left=col_left, top=col_top + Inches(0.62),
                         width=col_w, height=Inches(5.0), pt=14,
                         font_name=ds.font_body)


# ---- data_chart -----------------------------------------------------------
def _slide_data_chart(prs: Presentation, sd: dict, ds: DesignSystem) -> None:
    """
    Prompt schema: chart.{chart_type, title, description, x_axis_label,
    y_axis_label, data_series[]}, key_insight.
    Renders a real pptx chart when data is numeric, otherwise falls back to text.
    """
    slide = _blank_slide(prs)
    dn = sd.get("design_notes", {})
    _bg_solid(slide, dn.get("background", ds.bg_color))

    _title_box(slide, sd.get("title", ""), ds, top=Inches(0.3), pt=CONTENT_TITLE_PT)

    subtitle = sd.get("subtitle", "")
    if subtitle:
        _subtitle_box(slide, subtitle, ds, top=Inches(1.0))

    chart_spec  = sd.get("chart", {})
    data_series = chart_spec.get("data_series", [])
    chart_type_str = chart_spec.get("chart_type", "bar").lower()
    xl_chart_type  = CHART_TYPE_MAP.get(chart_type_str, XL_CHART_TYPE.BAR_CLUSTERED)

    chart_top = Inches(1.55) if subtitle else Inches(1.2)
    chart_h   = Inches(3.8)
    chart_left = Inches(0.8)
    chart_w    = Inches(8.4)

    rendered_chart = False
    if data_series:
        try:
            chart_data = CategoryChartData()
            first_series = data_series[0]
            labels = first_series.get("labels", [])
            chart_data.categories = labels

            for series in data_series:
                s_name = series.get("series_name", "Data")
                raw_values = series.get("values", [])
                numeric_vals = []
                for v in raw_values:
                    numeric_vals.append(float(str(v).replace(",", "").replace("$", "")))
                chart_data.add_series(s_name, numeric_vals)

            graphic_frame = slide.shapes.add_chart(
                xl_chart_type, chart_left, chart_top, chart_w, chart_h, chart_data
            )
            chart_obj = graphic_frame.chart
            chart_obj.has_legend = len(data_series) > 1

            x_label = chart_spec.get("x_axis_label", "")
            y_label = chart_spec.get("y_axis_label", "")
            if hasattr(chart_obj, "category_axis") and x_label:
                chart_obj.category_axis.has_title = True
                chart_obj.category_axis.axis_title.text_frame.paragraphs[0].text = x_label
            if hasattr(chart_obj, "value_axis") and y_label:
                chart_obj.value_axis.has_title = True
                chart_obj.value_axis.axis_title.text_frame.paragraphs[0].text = y_label

            # Apply series colour from JSON
            for idx, series in enumerate(data_series):
                s_color = series.get("color")
                if s_color and idx < len(chart_obj.series):
                    chart_obj.series[idx].format.fill.solid()
                    chart_obj.series[idx].format.fill.fore_color.rgb = _rgb(s_color)

            rendered_chart = True
            log.info("    Chart rendered: %s (%s)", chart_spec.get("title", ""), chart_type_str)
        except Exception as e:
            log.warning("    Chart rendering failed, using text fallback: %s", e)

    if not rendered_chart:
        lines = [f"[{chart_type_str.upper()} CHART] {chart_spec.get('title', '')}"]
        desc = chart_spec.get("description", "")
        if desc:
            lines.append(desc)
        for series in data_series:
            for lbl, val in zip(series.get("labels", []), series.get("values", [])):
                lines.append(f"  {lbl}: {val}")

        _, tf, _ = _textbox(slide, chart_left, chart_top, chart_w, chart_h)
        for i, line in enumerate(lines):
            para = tf.paragraphs[i] if i == 0 else tf.add_paragraph()
            para.text           = line
            para.font.size      = Pt(16 if i == 0 else 13)
            para.font.bold      = (i == 0)
            para.font.name      = ds.font_body
            para.font.color.rgb = _rgb(ds.primary_color if i == 0 else ds.text_on_light)
            para.alignment      = PP_ALIGN.CENTER if i == 0 else PP_ALIGN.LEFT

    # Key insight callout
    insight = sd.get("key_insight", {})
    if insight and insight.get("text"):
        _colored_box(
            slide, insight["text"],
            left=Inches(0.8), top=Inches(5.6),
            width=Inches(8.4), height=Inches(0.85),
            bg=insight.get("background_color", ds.primary_color),
            fg=insight.get("text_color", FB_WHITE),
            pt=14, font_name=ds.font_body,
        )


# ---- image_text -----------------------------------------------------------
def _slide_image_text(prs: Presentation, sd: dict, ds: DesignSystem) -> None:
    """Prompt schema: same as content but with an image area."""
    slide = _blank_slide(prs)
    dn = sd.get("design_notes", {})
    _bg_solid(slide, dn.get("background", ds.bg_color))

    _title_box(slide, sd.get("title", ""), ds, top=Inches(0.3), pt=CONTENT_TITLE_PT)

    image_desc = sd.get("image", {}).get("description", "")
    _, _, p = _textbox(slide, Inches(0.6), Inches(1.3), Inches(4.5), Inches(5.2))
    _fmt(p, f"[Image Placeholder]\n\n{image_desc}", 14,
         italic=True, color="#808080", font_name=ds.font_body)

    bullets = _get_bullets(sd)
    _render_bullets(slide, bullets, ds,
                    left=Inches(5.5), top=Inches(1.3),
                    width=Inches(4.0), height=Inches(5.2))


# ---- centered_content -----------------------------------------------------
def _slide_centered(prs: Presentation, sd: dict, ds: DesignSystem) -> None:
    slide = _blank_slide(prs)
    dn = sd.get("design_notes", {})
    _resolve_bg(slide, dn, ds)

    gc = dn.get("gradient_colors", [])
    bg = dn.get("background", ds.bg_color)
    is_dark = bool(gc) or (bg not in (FB_BG, "#FFFFFF", ds.bg_color))
    text_color = ds.text_on_dark if is_dark else ds.text_on_light

    accent_els = sd.get("accent_elements", {})
    line_color = accent_els.get("line_color", ds.accent_color)

    if accent_els.get("lines_before") in (True, "true"):
        _accent_line(slide, top=Inches(1.95),
                     left=Inches(1.2), width=Inches(7.6), color=line_color)

    _, _, p = _textbox(slide, Inches(0.8), Inches(2.15), Inches(8.4), Inches(1.5))
    _fmt(p, sd.get("title", ""), ds.title_pt,
         bold=True, color=text_color, align=PP_ALIGN.CENTER, font_name=ds.font_title)

    subtitle = sd.get("subtitle", "")
    if subtitle:
        _, _, p = _textbox(slide, Inches(0.8), Inches(3.75), Inches(8.4), Inches(0.7))
        _fmt(p, subtitle, 18, color=text_color, align=PP_ALIGN.CENTER,
             font_name=ds.font_body)

    points = sd.get("supporting_points", [])
    if points:
        _, tf, _ = _textbox(slide, Inches(1.5), Inches(4.55), Inches(7), Inches(1.8))
        tf.word_wrap = True
        for i, pt_text in enumerate(points):
            para = tf.paragraphs[i] if i == 0 else tf.add_paragraph()
            para.text           = f"\u2022 {pt_text}"
            para.font.size      = Pt(ds.body_pt)
            para.font.name      = ds.font_body
            para.font.color.rgb = _rgb(text_color)
            para.alignment      = PP_ALIGN.CENTER
            para.space_before   = Pt(4)

    if accent_els.get("lines_after") in (True, "true"):
        _accent_line(slide, top=Inches(6.45),
                     left=Inches(1.2), width=Inches(7.6), color=line_color)


# ---- comparison -----------------------------------------------------------
def _slide_comparison(prs: Presentation, sd: dict, ds: DesignSystem) -> None:
    """
    Prompt schema: left_column / right_column, each with
    label, background_color, text_color, items[].
    Renders as two side-by-side coloured boxes.
    """
    slide = _blank_slide(prs)
    dn = sd.get("design_notes", {})
    _bg_solid(slide, dn.get("background", ds.bg_color))
    _title_box(slide, sd.get("title", ""), ds, top=Inches(0.3), pt=CONTENT_TITLE_PT)

    col_width  = Inches(4.3)
    col_height = Inches(5.2)
    col_top    = Inches(1.3)

    col_defs = [
        (sd.get("left_column",  {}), Inches(0.5)),
        (sd.get("right_column", {}), Inches(5.3)),
    ]

    for col, col_left in col_defs:
        if not col:
            continue
        bg_c = col.get("background_color", ds.primary_color)
        fg_c = col.get("text_color", FB_WHITE)
        label = col.get("label", col.get("column_label", ""))

        box = slide.shapes.add_shape(1, col_left, col_top, col_width, col_height)
        box.fill.solid()
        box.fill.fore_color.rgb = _rgb(bg_c)
        box.line.color.rgb      = _rgb(bg_c)

        _, _, lp = _textbox(slide, col_left + Inches(0.1), col_top + Inches(0.1),
                            col_width - Inches(0.2), Inches(0.4))
        _fmt(lp, label, 16, bold=True, color=fg_c,
             align=PP_ALIGN.CENTER, font_name=ds.font_title)

        items = col.get("items", [])
        if items:
            _, tf, _ = _textbox(slide, col_left + Inches(0.15),
                                col_top + Inches(0.6),
                                col_width - Inches(0.3),
                                col_height - Inches(0.8))
            for i, item in enumerate(items):
                para = tf.paragraphs[i] if i == 0 else tf.add_paragraph()
                para.text           = f"\u2022 {item}"
                para.font.size      = Pt(12)
                para.font.name      = ds.font_body
                para.font.color.rgb = _rgb(fg_c)
                para.space_before   = Pt(3)
                para.space_after    = Pt(3)


# ---- closing --------------------------------------------------------------
def _slide_closing(prs: Presentation, sd: dict, ds: DesignSystem) -> None:
    slide = _blank_slide(prs)
    dn = sd.get("design_notes", {})
    _resolve_bg(slide, dn, ds, dark_default=ds.primary_color)

    _, _, p = _textbox(slide, MARGIN_L, Inches(1.5), CONTENT_W, Inches(1.3))
    _fmt(p, sd.get("title", ""), ds.title_pt,
         bold=True, color=ds.text_on_dark, align=PP_ALIGN.CENTER,
         font_name=ds.font_title)

    subtitle = sd.get("subtitle", "")
    if subtitle:
        _, _, p = _textbox(slide, MARGIN_L, Inches(2.9), CONTENT_W, Inches(0.7))
        _fmt(p, subtitle, 18,
             color=ds.text_on_dark, align=PP_ALIGN.CENTER, font_name=ds.font_body)

    next_steps = sd.get("next_steps", [])
    if next_steps:
        _render_str_list(slide, next_steps, ds.text_on_dark,
                         left=Inches(1.2), top=Inches(3.7),
                         width=Inches(7.6), height=Inches(1.0),
                         pt=14, font_name=ds.font_body,
                         bullet_prefix="\u2192 ")

    contact = sd.get("contact_info", {})
    lines   = [v for v in contact.values() if v]
    if lines:
        _, tf, _ = _textbox(slide, Inches(1), Inches(4.8), Inches(8), Inches(1.2))
        tf.word_wrap = True
        for i, line in enumerate(lines):
            para = tf.paragraphs[i] if i == 0 else tf.add_paragraph()
            para.text           = line
            para.font.size      = Pt(13)
            para.font.name      = ds.font_body
            para.font.color.rgb = _rgb(ds.text_on_dark)
            para.alignment      = PP_ALIGN.CENTER
            para.space_before   = Pt(3)

    cta = sd.get("cta_button", {})
    if cta and cta.get("text"):
        _colored_box(
            slide, cta["text"],
            left=Inches(2.8), top=Inches(6.25),
            width=Inches(4.4), height=Inches(0.65),
            bg=cta.get("background_color", ds.accent_color),
            fg=cta.get("text_color", FB_WHITE),
            pt=16, font_name=ds.font_body,
        )


# ---------------------------------------------------------------------------
# Dispatch table
# ---------------------------------------------------------------------------
HANDLERS: Dict[str, callable] = {
    "title":            _slide_title,
    "content":          _slide_content,
    "two_column":       _slide_two_column,
    "data_chart":       _slide_data_chart,
    "image_text":       _slide_image_text,
    "centered_content": _slide_centered,
    "comparison":       _slide_comparison,
    "closing":          _slide_closing,
}


# ---------------------------------------------------------------------------
# JSON loader
# ---------------------------------------------------------------------------
def load_json(path: str) -> dict:
    text = Path(path).read_text(encoding="utf-8").strip()
    if text.startswith("```"):
        text = re.sub(r"^```(?:json)?\n?", "", text)
        text = re.sub(r"\n?```$",          "", text)
    data = json.loads(text)
    log.info("Loaded JSON: %s", path)
    return data


# ---------------------------------------------------------------------------
# Core builder
# ---------------------------------------------------------------------------
def create_presentation(data: dict, output: str) -> str:
    prs = Presentation()
    prs.slide_width  = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    meta   = data.get("presentation_metadata", {})
    ds     = DesignSystem.from_dict(data.get("design_system_reference", {}))
    slides = data.get("slides", [])

    log.info("Building: %s  (%d slides)", meta.get("title", "Untitled"), len(slides))

    for sd in slides:
        snum  = sd.get("slide_number", "?")
        stype = sd.get("slide_type", "content")
        fn    = HANDLERS.get(stype)
        if fn is None:
            log.warning("Slide %s: unknown type %r — skipped", snum, stype)
            continue
        log.info("  Slide [%s] type=%s", snum, stype)
        fn(prs, sd, ds)
        _footer(prs.slides[-1], ds)
        _add_speaker_notes(prs.slides[-1], sd)

    Path(output).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)
    log.info("Saved → %s", output)
    return output


# ---------------------------------------------------------------------------
# Public entry point
# ---------------------------------------------------------------------------
def generate_presentation(
    json_file: Optional[str]        = None,
    output_file: Optional[str]      = None,
    prospect_company: Optional[str] = None,
) -> str:
    if json_file is None:
        if not prospect_company:
            raise ValueError("Provide json_file or prospect_company.")
        json_file = f"outputs/generated_slides/slides_{prospect_company}.json"

    if output_file is None:
        if not prospect_company:
            raise ValueError("Provide output_file or prospect_company.")
        slug = prospect_company.lower().replace(" ", "_")
        ts   = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_file = f"outputs/presentations/presentation_{slug}_{ts}.pptx"

    data   = load_json(json_file)
    result = create_presentation(data, output_file)
    return result


# ---------------------------------------------------------------------------
if __name__ == "__main__":
    generate_presentation(prospect_company="juniper")
